#!/usr/bin/env python3
"""Client-facing ICAI demo deck — business language, screenshot placeholders."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "ICAI_Publication_Portal_Client_Demo.pptx"
OUT_LEGACY = ROOT / "ICAI_Publication_Portal_Web_App_Demo.pptx"

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY_LIGHT = RGBColor(0x14, 0x32, 0x5C)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GOLD_SOFT = RGBColor(0xE8, 0xD4, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x47, 0x55, 0x69)
PLACEHOLDER_BG = RGBColor(0xF1, 0xF5, 0xF9)
PLACEHOLDER_BORDER = RGBColor(0xCB, 0xD5, 0xE1)

SLIDES = [
    {
        "layout": "title",
        "title": "ICAI Publication Portal",
        "subtitle": "Working Demo",
        "tagline": "Public Catalogue, Secure Reader, OTP Login & Admin CMS",
        "footer": "Prepared by Siyana Info Solutions Pvt. Ltd.",
        "notes": "Open with: this is a live working demonstration.",
        "placeholder": False,
    },
    {
        "title": "What We Demonstrated",
        "message": "A complete digital publication experience—from browse to read to manage.",
        "bullets": [
            "Public publication catalogue",
            "Login-based reading access",
            "Flipbook / PDF reader",
            "Admin upload and publishing control",
        ],
        "placeholder_label": "[ Optional: 4-panel overview graphic ]",
        "notes": "Set expectations for the next slides.",
    },
    {
        "title": "Public Catalogue View",
        "message": "Users discover ICAI publications before signing in to read.",
        "bullets": [
            "Browse ICAI publications in one catalogue",
            "Cover, title, committee, topic and synopsis visible",
            "Featured and latest releases supported",
            "Reading access is protected until login",
        ],
        "placeholder_label": "Screenshot:\nHomepage / publication listing",
        "notes": "Show filters and cards. File: 03-catalogue.png",
    },
    {
        "title": "Secure Login Flow",
        "message": "Sign-in options for visitors, members, and administrators.",
        "bullets": [
            "Non-members login using email OTP",
            "ICAI members connect via SSP SSO in production",
            "Admin users have separate CMS access",
            "Same reader access after successful authentication",
        ],
        "placeholder_label": "Screenshot:\nLogin page",
        "notes": "Demo OTP in controlled environment. File: 04-login.png",
    },
    {
        "title": "Flipbook Reader",
        "message": "Publications open in a secure, session-based reader.",
        "bullets": [
            "Opens only after login",
            "Page navigation, contents and search",
            "Download can be disabled or restricted",
            "Ideal for journals, handbooks and guides",
        ],
        "placeholder_label": "Screenshot:\nFlipbook / PDF reader",
        "notes": "Mention print/download disabled in demo. File: 05-flipbook.png",
    },
    {
        "title": "Admin CMS — Add Publication",
        "message": "Administrators add publications through a structured form.",
        "bullets": [
            "Upload PDF and cover image",
            "Add title, committee, topic, keywords and synopsis",
            "PDF, article, or PDF + article formats",
        ],
        "placeholder_label": "Screenshot:\nPublication upload form",
        "notes": "Walk top to bottom of form. File: 06-upload-form.png",
    },
    {
        "title": "Publishing & Access Control",
        "message": "Control how each publication appears and who can access it.",
        "bullets": [
            "Mark publication as featured",
            "Control download permission",
            "Set visibility for authenticated readers",
            "Draft, Published, Hidden, Archived workflow",
        ],
        "placeholder_label": "Screenshot:\nDistribution & status fields",
        "notes": "File: 07-publishing.png",
    },
    {
        "title": "Demo Publications Added",
        "message": "Sample ICAI records are ready for your review in the demo.",
        "bullets": [
            "The Chartered Accountant Journal — May 2026",
            "Handbook on Audit & Assurance 2026",
            "Covers and content attached; reader after login",
        ],
        "placeholder_label": "Screenshot:\nTwo publication cards",
        "notes": "Optional live open of handbook. Files: 08a, 08b",
    },
    {
        "title": "How It Works",
        "message": "Simple flows for readers and administrators.",
        "bullets": [
            "Readers: Catalogue → Login → Reader",
            "Admins: CMS → Upload → Publish → Catalogue",
        ],
        "placeholder_label": "Diagram:\nUser & Admin flows",
        "diagram": True,
        "notes": "Keep verbal—no technical stack on this slide.",
    },
    {
        "title": "Production Readiness",
        "message": "Clear steps from demonstration to ICAI production.",
        "bullets": [
            "ICAI SSP SSO integration",
            "ICAI-approved email gateway",
            "Secure hosting per ICAI policy",
            "Audit logs and analytics enhancement",
        ],
        "placeholder_label": "[ Roadmap / checklist graphic ]",
        "notes": "Position as partnership path, not gap list.",
    },
    {
        "title": "Demo Outcome",
        "message": "The portal is ready for ICAI-specific integration and rollout.",
        "bullets": [
            "Working portal flow demonstrated",
            "CMS publication management ready",
            "Secure reader access demonstrated",
            "Ready for ICAI-specific integration",
        ],
        "placeholder_label": "Screenshot:\nCollage (optional)",
        "notes": "Thank you; open Q&A.",
        "closing": True,
    },
]


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gold_bar(slide, top_in: float = 0) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(top_in), Inches(13.333), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def add_title_band(slide, title: str, message: str | None = None) -> None:
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.08), Inches(13.333), Inches(1.35)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8.5), Inches(0.65))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if message:
        mbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.82), Inches(8.5), Inches(0.45))
        mtf = mbox.text_frame
        mp = mtf.paragraphs[0]
        mp.text = message
        mp.font.size = Pt(13)
        mp.font.color.rgb = GOLD_SOFT


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(17)
        p.font.color.rgb = SLATE
        p.space_after = Pt(10)
        p.bullet = True


def add_screenshot_placeholder(
    slide, label: str, left: float, top: float, width: float, height: float, diagram: bool = False
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_BG
    shape.line.color.rgb = PLACEHOLDER_BORDER
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if diagram:
        p.text = (
            "Readers\n"
            "Catalogue  →  Login  →  Reader\n\n"
            "Administrators\n"
            "CMS  →  Upload  →  Publish  →  Catalogue"
        )
        p.font.size = Pt(14)
    else:
        p.text = label
        p.font.size = Pt(13)
    p.font.color.rgb = SLATE
    p.font.italic = True


def add_title_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)
    add_gold_bar(slide, 0)

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    p = tbox.text_frame.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.6))
    sp = sbox.text_frame.paragraphs[0]
    sp.text = data["subtitle"]
    sp.font.size = Pt(26)
    sp.font.color.rgb = GOLD
    sp.alignment = PP_ALIGN.CENTER

    tbox2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(10.8), Inches(0.5))
    tp = tbox2.text_frame.paragraphs[0]
    tp.text = data["tagline"]
    tp.font.size = Pt(16)
    tp.font.color.rgb = GOLD_SOFT
    tp.alignment = PP_ALIGN.CENTER

    fbox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4))
    fp = fbox.text_frame.paragraphs[0]
    fp.text = data["footer"]
    fp.font.size = Pt(12)
    fp.font.color.rgb = WHITE
    fp.alignment = PP_ALIGN.CENTER

  # screenshot placeholder bottom-right
    add_screenshot_placeholder(
        slide,
        "Screenshot:\nPortal header or catalogue (optional)",
        8.2,
        5.0,
        4.6,
        1.9,
    )

    slide.notes_slide.notes_text_frame.text = data.get("notes", "")


def add_content_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_gold_bar(slide, 0)
    add_title_band(slide, data["title"], data.get("message"))

    has_wide_placeholder = data.get("diagram") or len(data.get("bullets", [])) <= 3
    bullet_width = 5.8 if not has_wide_placeholder else 6.2
    add_bullets(slide, data["bullets"], 0.55, 1.65, bullet_width, 5.2)

    ph_left = 6.7 if bullet_width < 6 else 7.0
    add_screenshot_placeholder(
        slide,
        data.get("placeholder_label", "Screenshot"),
        ph_left,
        1.55,
        6.1,
        5.35,
        diagram=data.get("diagram", False),
    )

    if data.get("closing"):
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.95), Inches(13.333), Inches(0.55)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()
        t = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.35))
        t.text_frame.paragraphs[0].text = "Siyana Info Solutions Pvt. Ltd."
        t.text_frame.paragraphs[0].font.size = Pt(11)
        t.text_frame.paragraphs[0].font.color.rgb = WHITE
        t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    slide.notes_slide.notes_text_frame.text = data.get("notes", "")


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, data in enumerate(SLIDES):
        if data.get("layout") == "title":
            add_title_slide(prs, data)
        else:
            add_content_slide(prs, data)

    target = OUT
    try:
        prs.save(target)
    except PermissionError:
        target = OUT_LEGACY
        prs.save(target)
        print("Note: original .pptx was locked; saved alternate name.")
    print(f"Created: {target} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
